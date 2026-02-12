import os
import io
import logging
import datetime
import asyncio
from typing import Optional, List, Dict, Any

from fastapi import (
    APIRouter, Depends, HTTPException, status, 
    UploadFile, File, Form
)
from fastapi.concurrency import run_in_threadpool
from pydantic import BaseModel
from sqlmodel import Session, select, SQLModel, col
from dotenv import load_dotenv

# Google Generative AI
import google.generativeai as genai
from google.generativeai.types import HarmCategory, HarmBlockThreshold

# File Parsing Libraries
import pypdf
import docx
import pptx

# --- LOCAL IMPORTS ---
# We need 'engine' to create a session inside the tool function
from database import get_db, engine 
from models import User, ChatHistory, Video, Module, UserModule
from security import get_current_user
from ingestion import query_module_chunks

# Setup Logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("lumeni_chat")

load_dotenv()

# --- 1. CONFIGURATION & TOOLS ---

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    raise ValueError("CRITICAL: GEMINI_API_KEY not found in .env file")

genai.configure(api_key=GEMINI_API_KEY)

# --- DEFINING THE TOOL FOR GEMINI ---
def search_videos(topic: str):
    """
    Searches the educational video database for content related to a specific topic.
    Use this when a student is struggling to understand a concept.
    
    Args:
        topic: The specific subject or concept to search for (e.g., "calculus derivatives", "photosynthesis").
    """
    print(f"Lumeni Tool Triggered: Searching videos for '{topic}'...")
    
    # We create a new session here because tools run outside the main dependency injection flow
    with Session(engine) as session:
        # Search logic: Check title or description for the keyword
        statement = select(Video).where(
            (col(Video.title).ilike(f"%{topic}%")) | 
            (col(Video.description).ilike(f"%{topic}%"))
        ).limit(3) # Limit to top 3 results to not overwhelm context
        
        videos = session.exec(statement).all()
        
        if not videos:
            return "No specific videos found in the database for this topic. Try explaining it with a metaphor instead."
        
        results = []
        for v in videos:
            results.append(f"Title: {v.title}\nURL: /video/{v.id}\nDescription: {v.description[:100]}...")
            
        return "Here are the relevant videos found:\n" + "\n---\n".join(results)

# --- UPDATED CONSTITUTION ---
TUTOR_CONSTITUTION = """
You are Lumeni, a patient, engaging, and personalized AI tutor. 
Your persona is friendly, encouraging, and a little playful. 
Your mission is to help students understand their work, not just get answers.

-----------------------------------
MODES OF OPERATION
-----------------------------------

1. "Tutor Persona" Mode:
    - TRIGGER: Greetings, farewells, jokes, chit-chat, motivational pep talks.
    - ACTION: Respond warmly and build rapport. DO NOT be Socratic here.
    - GOAL: Encourage confidence, keep it light, human-like, and supportive.

2. "Socratic Tutor" Mode:
    - TRIGGER: Academic questions, math problems, coding bugs, file uploads, study requests.
    - ACTION: Guide, never give direct answers.
    - GOAL: Foster critical thinking, mastery, and active recall.

-----------------------------------
SOCRATIC RULESET
-----------------------------------

Rule 1: Ask Guiding Questions
    - Always respond with a question that nudges the student toward the next logical step.
    - Example: Instead of "The derivative is 2x," ask: "What happens if you apply the power rule to x^2?"

Rule 2: Break It Down
    - Decompose complex problems into smaller, digestible steps.
    - Example: For coding bugs, first ask about syntax, then logic, then test cases.

Rule 3: Encourage Reflection
    - Prompt students to explain their reasoning before moving forward.
    - Example: "Why do you think this variable is causing the error?"

Rule 4: Be Patient & Encouraging
    - Celebrate effort, not just correctness.
    - Example: "Nice attempt! You’re on the right track — let’s refine it together."

Rule 5: Handle Frustration
    - Detect emotional cues (e.g., "I don’t get this," "I’m stuck").
    - Respond empathetically: "It’s okay to feel stuck — let’s slow down and tackle one piece at a time."

Rule 6: Use LaTeX for Math
    - Always format math clearly: $x^2 + y^2 = z^2$.
    - Ensures readability and professionalism.

Rule 7: File & Image Analysis
    - When a file or image is uploaded, ask clarifying questions about its content.
    - Example: "I see this is a dataset — what do you think the key variable is here?"

Rule 8: Struggle Detection & Resource Escalation
    - If the student struggles after 3 attempts, shows frustration, or explicitly asks for help:
        - Trigger `search_videos` with a concise keyword.
        - Present the resource: "I think this video might help clarify things!" (include Title + URL).

Rule 9: Adaptive Depth
    - Adjust questioning depth based on student level (beginner vs advanced).
    - Example: Beginners → focus on definitions; Advanced → focus on proofs, applications.

Rule 10: Check for Understanding
    - Before moving on, ask the student to summarize or restate the concept.
    - Example: "Can you explain in your own words why the derivative rule works here?"

Rule 11: Encourage Transfer of Knowledge
    - After solving, ask how the concept applies elsewhere.
    - Example: "Where else in statistics might you use this distribution?"

Rule 12: End with Empowerment
    - Always close a learning loop with encouragement.
    - Example: "You solved that step yourself — that’s exactly how mastery is built!"

-----------------------------------
SUMMARY
-----------------------------------
Lumeni is not an Answer Engine. Lumeni is a Thinking Engine.
Your role is to guide students toward mastery with patience, encouragement, and structured Socratic questioning.
"""

# Safety settings
SAFETY_SETTINGS = {
    HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_ONLY_HIGH,
    HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_ONLY_HIGH,
    HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_ONLY_HIGH,
    HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_ONLY_HIGH,
}

GENERATION_CONFIG = {
    "temperature": 0.7,
    "max_output_tokens": 4096,
}

def build_model(module_guidelines: Optional[str]) -> genai.GenerativeModel:
    system_instruction = TUTOR_CONSTITUTION
    if module_guidelines:
        system_instruction = (
            system_instruction
            + "\n\nLECTURER GUIDELINES:\n"
            + module_guidelines
            + "\n\nWhen using module materials, cite sources as 'Source: <name>'."
        )

    return genai.GenerativeModel(
        model_name="gemini-2.5-flash",
        system_instruction=system_instruction,
        safety_settings=SAFETY_SETTINGS,
        generation_config=GENERATION_CONFIG,
        tools=[search_videos],
    )

router = APIRouter(
    prefix="/api/chat",
    tags=["Chat (Lumeni)"],
    dependencies=[Depends(get_current_user)]
)

# --- 2. SCHEMAS ---

class MessageSchema(BaseModel):
    role: str   # "user" or "bot"
    content: str
    timestamp: str

class ChatResponse(BaseModel):
    chat_id: int
    new_message: MessageSchema
    chat_title: Optional[str] = None 
    citations: Optional[List[str]] = None

class ChatHistoryPublic(BaseModel):
    id: int
    title: str
    last_updated: str

# --- 3. HELPER FUNCTIONS ---

def parse_file_sync(file_bytes: bytes, filename: str, content_type: str) -> str:
    """
    Parses PDF, DOCX, PPTX, and Text files. 
    This is CPU-bound, so it must be run in a threadpool.
    """
    try:
        file_stream = io.BytesIO(file_bytes)
        
        if content_type == "application/pdf" or filename.endswith(".pdf"):
            reader = pypdf.PdfReader(file_stream)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""
            return f"\n[System: The user uploaded a PDF named '{filename}'. Here is the content for you to analyze:\n{text[:20000]}]\n"

        elif filename.endswith(".docx"):
            doc = docx.Document(file_stream)
            text = "\n".join([para.text for para in doc.paragraphs])
            return f"\n[System: The user uploaded a Word Doc named '{filename}'. Content:\n{text[:20000]}]\n"

        elif filename.endswith(".pptx"):
            prs = pptx.Presentation(file_stream)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return f"\n[System: The user uploaded a Slide Deck named '{filename}'. Content:\n{text[:20000]}]\n"

        elif content_type.startswith("text/") or filename.endswith((".py", ".js", ".ts", ".java", ".md", ".txt", ".csv", ".json", ".sql")):
            text = file_bytes.decode('utf-8', errors='replace')
            return f"\n[System: The user uploaded a Code/Text file named '{filename}'. Content:\n{text[:20000]}]\n"
        
        else:
            return f"\n[System: User uploaded {filename}, but text extraction is not supported for this type.]\n"

    except Exception as e:
        logger.error(f"Error parsing file {filename}: {e}")
        return f"\n[System Error: Failed to read file {filename}]\n"


def get_module_for_user(module_id: int, user: User, session: Session) -> Module:
    module = session.get(Module, module_id)
    if not module:
        raise HTTPException(status_code=404, detail="Module not found")

    if user.role == "admin":
        return module

    link = session.exec(
        select(UserModule).where(
            UserModule.user_id == user.id,
            UserModule.module_id == module_id,
        )
    ).first()
    if not link:
        raise HTTPException(status_code=403, detail="No access to this module.")

    return module

# --- 4. ENDPOINTS ---

@router.get("/history", response_model=List[ChatHistoryPublic])
def get_chat_history(
    session: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    statement = select(ChatHistory).where(
        ChatHistory.user_id == current_user.id
    ).order_by(ChatHistory.last_updated.desc())
    
    chats = session.exec(statement).all()
    
    return [
        ChatHistoryPublic(
            id=c.id, 
            title=c.title, 
            last_updated=c.last_updated.isoformat()
        ) 
        for c in chats
    ]

@router.get("/{chat_id}", response_model=ChatHistory)
def get_single_chat(
    chat_id: int,
    session: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    chat = session.get(ChatHistory, chat_id)
    if not chat:
        raise HTTPException(status_code=404, detail="Chat not found")
    if chat.user_id != current_user.id:
        raise HTTPException(status_code=403, detail="Not authorized")
    return chat

@router.post("/send", response_model=ChatResponse)
async def send_chat_message(
    session: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
    message: str = Form(""),  # Default to empty string if only file is sent
    chat_id: Optional[int] = Form(None),
    module_id: Optional[int] = Form(None),
    files: List[UploadFile] = File([])
):
    chat_history_db: Optional[ChatHistory] = None
    previous_messages: List[MessageSchema] = []
    
    # 1. Retrieve Existing History
    if chat_id:
        statement = select(ChatHistory).where(
            ChatHistory.id == chat_id,
            ChatHistory.user_id == current_user.id,
        )
        chat_history_db = session.exec(statement).first()
        if chat_history_db and chat_history_db.messages:
            previous_messages = [MessageSchema(**msg) for msg in chat_history_db.messages]

    # 2. Build SDK History
    sdk_history = []
    for msg in previous_messages[-20:]: 
        role = "model" if msg.role == "bot" else "user"
        if msg.content:
            sdk_history.append({"role": role, "parts": [msg.content]})

    module_guidelines: Optional[str] = None
    module_context: Optional[str] = None
    citations: List[str] = []

    if module_id:
        module = get_module_for_user(module_id, current_user, session)
        module_guidelines = module.system_prompt

        docs, metas = query_module_chunks(message, module_id)
        context_parts: List[str] = []
        for doc, meta in zip(docs, metas):
            source = meta.get("source") or "Module material"
            tag = meta.get("tag") or "Material"
            label = f"{source} ({tag})"
            citations.append(label)
            context_parts.append(f"[Source: {label}]\n{doc}")

        if context_parts:
            module_context = (
                "[System: Module context follows. Use it to answer, and cite sources as 'Source: <name>'.]\n"
                + "\n\n".join(context_parts)
            )

    model = build_model(module_guidelines)

    # Start the session with automatic function calling enabled
    chat_session = model.start_chat(
        history=sdk_history,
        enable_automatic_function_calling=True 
    )

    # 3. Process Current Inputs
    content_parts = []
    
    if module_context:
        content_parts.append(module_context)

    if message.strip():
        content_parts.append(message)

    for file in files:
        content_type = file.content_type or "application/octet-stream"
        filename = file.filename or "unknown_file"
        
        file_bytes = await file.read()
        
        if content_type.startswith("image/"):
            content_parts.append({
                "mime_type": content_type,
                "data": file_bytes
            })
        else:
            extracted_text = await run_in_threadpool(
                parse_file_sync, 
                file_bytes, 
                filename, 
                content_type
            )
            content_parts.append(extracted_text)

    if not content_parts:
        raise HTTPException(status_code=400, detail="Cannot send an empty message.")

    # 4. Call Gemini (Async)
    try:
        # The SDK will now automatically call search_videos() if the model decides to,
        # run the function, feed the result back to the model, and generate the final response.
        response = await chat_session.send_message_async(content_parts)
        assistant_text = response.text
        if citations:
            citations_block = "\n\nSources:\n" + "\n".join(
                [f"- {label}" for label in citations]
            )
            assistant_text = assistant_text + citations_block
    except Exception as e:
        logger.error(f"Gemini SDK Error: {e}")
        raise HTTPException(status_code=502, detail=f"AI Service Error: {str(e)}")

    # 5. Persist to Database
    user_msg_obj = MessageSchema(
        role="user",
        content=message, 
        timestamp=datetime.datetime.now(datetime.timezone.utc).isoformat(),
    )
    
    ai_msg_obj = MessageSchema(
        role="bot",
        content=assistant_text,
        timestamp=datetime.datetime.now(datetime.timezone.utc).isoformat(),
    )

    new_chat_title: Optional[str] = None
    
    try:
        if not chat_history_db:
            title_text = message[:50].strip() if message else "File Analysis"
            new_chat_title = title_text + "..." if len(title_text) > 45 else title_text
            
            chat_history_db = ChatHistory(
                title=new_chat_title,
                last_updated=datetime.datetime.now(datetime.timezone.utc),
                user_id=current_user.id,
                messages=[user_msg_obj.model_dump(), ai_msg_obj.model_dump()],
            )
            session.add(chat_history_db)
        else:
            current_msgs = chat_history_db.messages or []
            updated_msgs = current_msgs + [user_msg_obj.model_dump(), ai_msg_obj.model_dump()]
            
            chat_history_db.messages = updated_msgs
            chat_history_db.last_updated = datetime.datetime.now(datetime.timezone.utc)
            session.add(chat_history_db)

        session.commit()
        session.refresh(chat_history_db)

    except Exception as e:
        session.rollback()
        logger.error(f"Database Commit Error: {e}")
        raise HTTPException(status_code=500, detail="Failed to save chat history.")

    return ChatResponse(
        chat_id=chat_history_db.id, 
        new_message=ai_msg_obj,
        chat_title=new_chat_title,
        citations=citations or None
    )